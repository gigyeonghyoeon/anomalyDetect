"""Generate assignment presentation (PPTX) from experiment results."""

from __future__ import annotations

import json
from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
RESULTS = ROOT / "결과모음" / "results"
OUT_DIR = ROOT / "결과모음"
CHART_PATH = OUT_DIR / "ppt_exp2_auroc_chart.png"
PPTX_PATH = OUT_DIR / "발표_흉부Xray_이상탐지.pptx"

# Colors
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x0E, 0x7C, 0x86)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)


def load_data() -> tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame, dict]:
    p1 = pd.read_csv(RESULTS / "phase1_model_comparison.csv")
    p2 = pd.read_csv(RESULTS / "phase2_preprocess_comparison.csv")
    p3 = pd.read_csv(RESULTS / "hyperparam_results.csv")
    with open(RESULTS / "best_final_config.json", encoding="utf-8") as f:
        best = json.load(f)
    return p1, p2, p3, best


def make_chart(best: dict) -> Path:
    plt.rcParams["font.family"] = ["Malgun Gothic", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False

    phases = ["Phase 1\n(basic)", "Phase 2\n(enhanced)", "Phase 3\n(HP tune)"]
    values = [0.559, 0.619, best["exp2_auroc"]]
    colors = ["#94a3b8", "#0e7490", "#0f766e"]

    fig, ax = plt.subplots(figsize=(7, 4))
    bars = ax.bar(phases, values, color=colors, width=0.55, edgecolor="white")
    ax.set_ylim(0, 0.75)
    ax.set_ylabel("Exp2 AUROC (RSNA)", fontsize=11)
    ax.set_title("단계별 RSNA 일반화 성능", fontsize=13, fontweight="bold")
    ax.axhline(0.5, color="#cbd5e1", linestyle="--", linewidth=1, label="random (0.5)")
    for bar, val in zip(bars, values):
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            bar.get_height() + 0.012,
            f"{val:.3f}",
            ha="center",
            va="bottom",
            fontsize=11,
            fontweight="bold",
        )
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.tight_layout()
    fig.savefig(CHART_PATH, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return CHART_PATH


def set_slide_bg(slide, color: RGBColor = LIGHT_BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "흉부 X-ray 기반\n폐렴 이상탐지 실험"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    sub = tf.add_paragraph()
    sub.text = "비지도 학습 · PatchCore · 3단계 실험"
    sub.font.size = Pt(20)
    sub.font.color.rgb = RGBColor(0xA5, 0xF3, 0xFC)
    sub.space_before = Pt(16)

    info = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(8), Inches(1))
    itf = info.text_frame
    ip = itf.paragraphs[0]
    ip.text = "(과목명)  |  (학번)  |  (이름)  |  2026.06"
    ip.font.size = Pt(14)
    ip.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)


def add_section_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, ACCENT)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.4), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sp = tf.add_paragraph()
        sp.text = subtitle
        sp.font.size = Pt(18)
        sp.font.color.rgb = RGBColor(0xCC, 0xFB, 0xF1)
        sp.space_before = Pt(12)


def add_bullet_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    note: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.8))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.15), Inches(1.2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.8), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.space_after = Pt(10)
        p.level = 0

    if note:
        nb = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(8.8), Inches(0.6))
        np = nb.text_frame.paragraphs[0]
        np.text = note
        np.font.size = Pt(12)
        np.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        np.font.italic = True


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    footnote: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.8))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.3), Inches(9), Inches(0.45 * n_rows)).table

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    if footnote:
        fb = slide.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(8.8), Inches(0.5))
        fp = fb.text_frame.paragraphs[0]
        fp.text = footnote
        fp.font.size = Pt(11)
        fp.font.color.rgb = GRAY


def add_chart_slide(prs: Presentation, title: str, chart_path: Path, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    slide.shapes.add_picture(str(chart_path), Inches(0.5), Inches(1.2), width=Inches(5.2))

    body = slide.shapes.add_textbox(Inches(5.9), Inches(1.5), Inches(3.8), Inches(4.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.space_after = Pt(8)


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_items: list[str],
    right_title: str,
    right_items: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    for x, col_title, items in [
        (0.6, left_title, left_items),
        (5.1, right_title, right_items),
    ]:
        box = slide.shapes.add_textbox(Inches(x), Inches(1.2), Inches(4.2), Inches(5))
        tf = box.text_frame
        tf.word_wrap = True
        hp = tf.paragraphs[0]
        hp.text = col_title
        hp.font.size = Pt(18)
        hp.font.bold = True
        hp.font.color.rgb = ACCENT
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(15)
            p.font.color.rgb = GRAY
            p.space_before = Pt(6)


def add_closing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8.4), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "감사합니다"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    sp = tf.add_paragraph()
    sp.text = "Q & A"
    sp.font.size = Pt(24)
    sp.font.color.rgb = RGBColor(0xA5, 0xF3, 0xFC)
    sp.alignment = PP_ALIGN.CENTER
    sp.space_before = Pt(20)


def build_presentation() -> Path:
    p1, p2, p3, best = load_data()
    chart_path = make_chart(best)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_bullet_slide(
        prs,
        "연구 배경 & 문제",
        [
            "폐렴 조기 탐지는 임상적으로 중요하지만, 질환 데이터 라벨링 비용이 큼",
            "정상(NORMAL) 영상만으로 학습 → 이상탐지 (Unsupervised AD)",
            "핵심 과제: 학습 도메인과 다른 RSNA 데이터에서도 잘 동작하는가?",
        ],
    )

    add_bullet_slide(
        prs,
        "연구 목표",
        [
            "Chest X-ray NORMAL만으로 Conv AE / U-Net AE / PatchCore 학습",
            "Exp1 (Chest X-ray test) vs Exp2 (RSNA) 성능 비교",
            "전처리·HP 튜닝으로 Exp2 AUROC 최대화",
            "우승 기준: Exp2 (RSNA) AUROC 최대",
        ],
        note="평가 지표: AUROC (주), F1 / Precision / Recall (보조)",
    )

    add_table_slide(
        prs,
        "데이터 & 실험 설계",
        ["구분", "데이터", "역할"],
        [
            ["학습", "Chest X-ray NORMAL", "비지도 학습 (정상만)"],
            ["Exp1", "Chest X-ray test", "동일 도메인 검증"],
            ["Exp2", "RSNA Pneumonia", "외부 도메인 · 일반화"],
        ],
        footnote="환경: Python · PyTorch · AWS EC2 GPU  |  3 Phase 실험 (총 21 runs)",
    )

    add_section_slide(prs, "방법론", "모델 · 전처리 · 평가")

    add_two_column_slide(
        prs,
        "모델 & 전처리",
        "3종 모델",
        [
            "Conv AE — 재구성 오차 기반",
            "U-Net AE — skip connection AE",
            "PatchCore — ResNet18 + memory bank, k-NN 거리",
        ],
        "전처리 (basic / enhanced)",
        [
            "basic: grayscale → resize → normalize",
            "enhanced: + lung crop + CLAHE",
            "domain_gap = Exp1 AUROC − Exp2 AUROC",
        ],
    )

    add_table_slide(
        prs,
        "Phase 1 — 모델 비교 (basic)",
        ["모델", "Exp1 AUROC", "Exp2 AUROC", "판정"],
        [
            ["Conv AE", "0.486", "0.478", "≈ random"],
            ["U-Net AE", "0.493", "0.490", "≈ random"],
            ["PatchCore", "0.716", "0.559", "★ 우승"],
        ],
        footnote="AE 계열 AUROC ~0.48 → 재구성 오차만으로 RSNA 일반화 불가",
    )

    add_table_slide(
        prs,
        "Phase 2 — 전처리 비교 (PatchCore)",
        ["전처리", "Exp1", "Exp2", "domain_gap"],
        [
            ["basic", "0.716", "0.559", "0.157"],
            ["enhanced", "0.663", "0.619", "0.044"],
        ],
        footnote="enhanced: Exp2 +6.0%p ↑, domain_gap 대폭 감소 → RSNA 일반화 개선",
    )

    top5 = (
        p3[p3["experiment"] == "exp2"]
        .sort_values("auroc", ascending=False)
        .drop_duplicates("run_id")
        .head(5)
    )
    rows = []
    for _, r in top5.iterrows():
        exp1 = p3[(p3["run_id"] == r["run_id"]) & (p3["experiment"] == "exp1")]["auroc"].iloc[0]
        rows.append([r["run_id"], f"{exp1:.3f}", f"{r['auroc']:.3f}"])
    add_table_slide(
        prs,
        "Phase 3 — HP 탐색 Top 5 (Exp2 AUROC)",
        ["run_id", "Exp1", "Exp2"],
        rows,
        footnote="16조합: lr × image(224/256) × k(5/9) × coreset(0.05/0.1)  |  img256·k5·cs01 최고",
    )

    add_chart_slide(
        prs,
        "단계별 Exp2 AUROC 개선",
        chart_path,
        [
            "Phase 1 → 3: +7.8%p",
            "0.559 → 0.619 → 0.637",
            "enhanced 전처리 효과 큼",
            "HP 튜닝으로 추가 +1.8%p",
        ],
    )

    final_row = p3[(p3["run_id"] == best["run_id"]) & (p3["experiment"] == "exp2")].iloc[0]
    exp1_row = p3[(p3["run_id"] == best["run_id"]) & (p3["experiment"] == "exp1")].iloc[0]
    add_table_slide(
        prs,
        "최종 선정 모델",
        ["항목", "Exp1", "Exp2"],
        [
            ["AUROC", f"{exp1_row['auroc']:.3f}", f"{final_row['auroc']:.3f}"],
            ["F1", f"{exp1_row['f1']:.3f}", f"{final_row['f1']:.3f}"],
            ["Precision", f"{exp1_row['precision']:.3f}", f"{final_row['precision']:.3f}"],
            ["Recall", f"{exp1_row['recall']:.3f}", f"{final_row['recall']:.3f}"],
        ],
        footnote=(
            f"PatchCore + enhanced | lr1e-3 | img256 | k5 | cs0.1  "
            f"| Exp2 Recall {final_row['recall']:.2f} (민감↑, Precision {final_row['precision']:.2f})"
        ),
    )

    add_bullet_slide(
        prs,
        "결론 & 한계",
        [
            "PatchCore만 실용적 — AE는 AUROC ~0.48 (분류 불가)",
            "enhanced 전처리 + HP 튜닝 → Exp2 AUROC 0.637 (최종)",
            "domain_gap 0.157 → 0.052로 완화, RSNA 일반화 개선",
            "한계: AUROC 0.64 수준, 오탐(precision) 높음 → 임상 적용엔 추가 연구 필요",
        ],
    )

    add_closing_slide(prs)

    prs.save(PPTX_PATH)
    return PPTX_PATH


if __name__ == "__main__":
    path = build_presentation()
    print(f"Presentation saved to {path}")
